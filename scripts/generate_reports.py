# -*- coding: utf-8 -*-
import json, os
from reportlab.lib.pagesizes import landscape, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

DATA_FILE = "data/input.json"
OUT_DIR = "output"
PDF_PATH = os.path.join(OUT_DIR, "ВОЛС_руководителю.pdf")
PPTX_PATH = os.path.join(OUT_DIR, "ВОЛС_руководителю.pptx")

def ensure_fonts():
    # Нужны кириллические шрифты для PDF
    djv = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
    djv_b = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
    pdfmetrics.registerFont(TTFont("DejaVuSans", djv))
    pdfmetrics.registerFont(TTFont("DejaVuSans-Bold", djv_b))

def build_pdf(d):
    os.makedirs(OUT_DIR, exist_ok=True)
    ensure_fonts()

    styles = getSampleStyleSheet()
    styles["Normal"].fontName = "DejaVuSans"
    styles["Normal"].fontSize = 12
    styles["Heading1"].fontName = "DejaVuSans-Bold"
    styles["Heading2"].fontName = "DejaVuSans-Bold"

    doc = SimpleDocTemplate(PDF_PATH, pagesize=landscape(A4))
    E = []
    E.append(Paragraph(f"Информация по бездоговорному подвесу ВОЛС (на {d['as_of']})", styles["Heading1"]))
    E.append(Spacer(1, 18))

    # 1. Общая ситуация
    E.append(Paragraph("1. Общая ситуация", styles["Heading2"]))
    t1 = [
        ["Показатель", "Значение"],
        ["Выявлено бездоговорных опор", f"{d['totals']['found']:,}".replace(',', ' ')],
        ["Узаконено", f"{d['totals']['legalized']:,}".replace(',', ' ')],
        ["Демонтировано в 2025", f"{d['totals']['removed_2025']:,}".replace(',', ' ') + f" (в 2024 — {d['totals']['removed_2024']:,}".replace(',', ' ') + ")"],
        ["В работе", f"{d['totals']['in_work']:,}".replace(',', ' ') + f" (ПАО «Ростелеком» — {d['totals']['rostelecom']:,}".replace(',', ' ') + ")"],
    ]
    tbl1 = Table(t1, colWidths=[300, 300])
    tbl1.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.lightgrey),
        ("GRID", (0,0), (-1,-1), 0.5, colors.grey),
        ("FONTNAME", (0,0), (-1,-1), "DejaVuSans"),
        ("FONTSIZE", (0,0), (-1,-1), 12),
    ]))
    E.append(tbl1); E.append(Spacer(1, 12))

    # 2. Демонтаж 2025
    E.append(Paragraph("2. Демонтаж 2025", styles["Heading2"]))
    t2 = [["Филиал", "Демонтировано / Примечание"]]
    for r in d["dismantled_2025"]:
        t2.append([r["branch"], f"{r['count']:,}".replace(',', ' ') + (f" ({r['notes']})" if r.get("notes") else "")])
    tbl2 = Table(t2, colWidths=[300, 300])
    tbl2.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.lightgrey),
        ("GRID", (0,0), (-1,-1), 0.5, colors.grey),
        ("FONTNAME", (0,0), (-1,-1), "DejaVuSans"),
        ("FONTSIZE", (0,0), (-1,-1), 12),
    ]))
    E.append(tbl2); E.append(Spacer(1, 12))

    # 3. Ключевые филиалы с цветовой маркировкой риска
    E.append(Paragraph("3. Ключевые филиалы (зоны риска)", styles["Heading2"]))
    t3 = [["Филиал", "Опоры в работе / Примечание", "Риск"]]
    risk_color = {"high": colors.red, "medium": colors.orange, "low": colors.darkgoldenrod}
    for r in d["key_branches"]:
        t3.append([r["branch"], f"{r['in_work']} ({r['note']})", r["risk"]])
    tbl3 = Table(t3, colWidths=[250, 300, 50])
    tbl3_style = [
        ("BACKGROUND", (0,0), (-1,0), colors.lightgrey),
        ("GRID", (0,0), (-1,-1), 0.5, colors.grey),
        ("FONTNAME", (0,0), (-1,-1), "DejaVuSans"),
        ("FONTSIZE", (0,0), (-1,-1), 12),
        ("ALIGN", (2,1), (2,-1), "CENTER"),
    ]
    # раскрасим колонку «Риск»
    for i, r in enumerate(d["key_branches"], start=1):
        tbl3_style.append(("TEXTCOLOR", (2, i), (2, i), risk_color[r["risk"]]))
    tbl3.setStyle(TableStyle(tbl3_style))
    E.append(tbl3); E.append(Spacer(1, 12))

    # 4. Ростелеком
    E.append(Paragraph("4. ПАО «Ростелеком» (общий объём: " + f"{d['totals']['rostelecom']:,}".replace(',', ' ') + " опор)", styles["Heading2"]))
    t4 = [["Филиал", "Статус"]]
    for r in d["rostelecom"]:
        t4.append([r["branch"], r["note"]])
    tbl4 = Table(t4, colWidths=[250, 350])
    tbl4.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.lightgrey),
        ("GRID", (0,0), (-1,-1), 0.5, colors.grey),
        ("FONTNAME", (0,0), (-1,-1), "DejaVuSans"),
        ("FONTSIZE", (0,0), (-1,-1), 12),
    ]))
    E.append(tbl4); E.append(Spacer(1, 12))

    # 5. Вывод
    E.append(Paragraph("📌 Вывод для руководителя", styles["Heading2"]))
    E.append(Paragraph("— Темпы демонтажа просели (2025 vs 2024).", styles["Normal"]))
    E.append(Paragraph("— «Ростелеком» — основной проблемный контрагент.", styles["Normal"]))
    E.append(Paragraph("— ЮЗ ЭС, Краснодарские ЭС, Армавирские ЭС формируют львиную долю риска.", styles["Normal"]))
    E.append(Paragraph("— Нужны: ускорение подписания, усиление демонтажа, персональная ответственность директорів филиалов.", styles["Normal"]))

    doc.build(E)

def build_pptx(d):
    os.makedirs(OUT_DIR, exist_ok=True)
    prs = Presentation()

    # Слайд 1 — заголовок
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Информация по бездоговорному подвесу ВОЛС"
    slide.placeholders[1].text = f"АО «Россети Кубань» • {d['as_of']}"

    # Слайд 2 — Общая ситуация
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. Общая ситуация"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for line in [
        f"Выявлено: {d['totals']['found']:,}".replace(',', ' '),
        f"Узаконено: {d['totals']['legalized']:,}".replace(',', ' '),
        f"Демонтировано 2025: {d['totals']['removed_2025']:,}".replace(',', ' ') + f" (2024: {d['totals']['removed_2024']:,}".replace(',', ' ') + ")",
        f"В работе: {d['totals']['in_work']:,}".replace(',', ' ') + f" (Ростелеком — {d['totals']['rostelecom']:,}".replace(',', ' ') + ")"
    ]:
        p = tf.add_paragraph(); p.text = line; p.font.size = Pt(20)

    # Слайд 3 — Демонтаж 2025
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2. Демонтаж 2025"
    tf = slide.placeholders[1].text_frame; tf.clear()
    for r in d["dismantled_2025"]:
        p = tf.add_paragraph()
        p.text = f"{r['branch']} — {r['count']:,}".replace(',', ' ') + (f" ({r['notes']})" if r.get("notes") else "")
        p.font.size = Pt(20)

    # Слайд 4 — Ключевые филиалы (раскраска)
    colors_map = {"high": RGBColor(0xFF, 0x00, 0x00),
                  "medium": RGBColor(0xFF, 0xA5, 0x00),
                  "low": RGBColor(0xFF, 0xD7, 0x00)}
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3. Ключевые филиалы — зоны риска"
    tf = slide.placeholders[1].text_frame; tf.clear()
    for r in d["key_branches"]:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"{r['branch']} — {r['in_work']} ({r['note']})"
        run.font.size = Pt(20)
        run.font.color.rgb = colors_map[r["risk"]]

    # Слайд 5 — Ростелеком
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "4. ПАО «Ростелеком»"
    tf = slide.placeholders[1].text_frame; tf.clear()
    p0 = tf.add_paragraph()
    p0.text = f"Общий объём: {d['totals']['rostelecom']:,}".replace(',', ' ') + " опор"
    p0.font.size = Pt(20)
    for r in d["rostelecom"]:
        p = tf.add_paragraph()
        p.text = f"{r['branch']} — {r['note']}"
        p.font.size = Pt(20)

    # Слайд 6 — Вывод
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "📌 Вывод для руководителя"
    tf = slide.placeholders[1].text_frame; tf.clear()
    for line in [
        "Темпы демонтажа просели (2025 vs 2024).",
        "«Ростелеком» — основной проблемный контрагент.",
        "ЮЗ ЭС, Краснодарские ЭС, Армавирские ЭС — основная зона риска.",
        "Нужны: ускорение подписания, усиление демонтажа, персональная ответственность."
    ]:
        p = tf.add_paragraph(); p.text = line; p.font.size = Pt(20)

    prs.save(PPTX_PATH)

def main():
    with open(DATA_FILE, "r", encoding="utf-8") as f:
        data = json.load(f)
    build_pdf(data)
    build_pptx(data)
    print("Готово:", PDF_PATH, PPTX_PATH)

if __name__ == "__main__":
    main()
